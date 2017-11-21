# Thomas Aleks Miller
# My professor doesn't give us his ppt for his lectures in Biochemistry and I take photos
# of each slide and make my own powerpoint. Saves the time of drag and dropping each photo
# as well as formatting each photo.
import pptx
import pptx.util
import glob
import scipy.misc
import os
import os.path

# Name your presentation
OUTPUT_TAG = 'Topic_69'
# Initialize Powerpoint Presentation
prs = pptx.Presentation()

# Store location of photos in a variable
photoloc = '/pptphotos'

# Initialize slide layout for presentation
blank_slide_layout = prs.slide_layouts[6]

# Initialize an empty list to store filenames of photos
photolist = []

# for every file in the folder add to photolist
for root, dirs, filenames in os.walk(photoloc):
    for f in filenames:
	
		photolist.append(f)

# initialize photo list to put each photos location
photos = []
slicepl = photolist[1:]
# for every file name add it's directory
for i in slicepl:

	a = "pptphotos/" + i

	photos.append(a)
	

# for every photo create a new slide and add one photo per slide
for p in photos:

	slide = prs.slides.add_slide(blank_slide_layout)


	pic = slide.shapes.add_picture(p,pptx.util.Inches(0.5), pptx.util.Inches(1.75), width=pptx.util.Inches(9), height=pptx.util.Inches(5))

prs.save("%s.pptx" % OUTPUT_TAG)
